from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT = "/Users/kusumathatavarthi/jira_ai_chatbot_artifacts/Jira_AI_Copilot_5min_Demo.pptx"

prs = Presentation()


def apply_dark_theme(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(14, 18, 30)

    accent = slide.shapes.add_shape(
        autoshape_type_id=1,
        left=Inches(11.1),
        top=Inches(0.0),
        width=Inches(0.8),
        height=Inches(1.0),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(186, 21, 21)
    accent.line.fill.background()


def style_title(shape, text):
    shape.text = text
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.font.size = Pt(42)
    p.font.bold = False
    p.font.color.rgb = RGBColor(240, 242, 248)


def add_subtitle_box(slide, text):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.0), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(21)
    p.font.color.rgb = RGBColor(190, 198, 216)


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_theme(slide)
    style_title(slide.shapes.title, title)
    add_subtitle_box(slide, subtitle)


def add_bullets_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_theme(slide)
    style_title(slide.shapes.title, title)

    content = slide.shapes.add_textbox(Inches(0.85), Inches(1.8), Inches(11.2), Inches(4.8))
    tf = content.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {b}"
        p.level = 0
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(227, 233, 246)


def add_demo_flow_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_theme(slide)
    style_title(slide.shapes.title, "Live Demo Flow (2 minutes)")

    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(12.0)
    height = Inches(4.5)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    steps = [
        "1) Ask: Any P1?",
        "2) Ask: Who is working on P1 tickets?",
        "3) Ask: Why was Sprint 24 delayed? executive summary",
        "4) Show citations + index used (trust and traceability)",
        "5) Show P1 owner summary for instant action",
    ]
    for i, step in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = step
        p.font.size = Pt(26)
        p.font.color.rgb = RGBColor(230, 235, 248)

    hint = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5))
    htf = hint.text_frame
    htf.text = "Red-highlight focus: P1 ownership, risk, and citations"
    htf.paragraphs[0].font.size = Pt(16)
    htf.paragraphs[0].font.color.rgb = RGBColor(229, 67, 67)


def add_highlight_slide(title, subtitle, points):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_theme(slide)
    style_title(slide.shapes.title, title)

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(12), Inches(0.6))
    stf = subtitle_box.text_frame
    stf.text = subtitle
    stf.paragraphs[0].font.size = Pt(20)
    stf.paragraphs[0].font.color.rgb = RGBColor(229, 67, 67)

    body = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(12), Inches(4.5))
    btf = body.text_frame
    for i, point in enumerate(points):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = f"▸  {point}"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(225, 231, 245)


def add_text_panel_slide(title, lines, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    top = Inches(1.3)
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(12), Inches(0.5))
        stf = subtitle_box.text_frame
        stf.text = subtitle
        stf.paragraphs[0].font.size = Pt(18)
        stf.paragraphs[0].font.color.rgb = RGBColor(60, 90, 160)
        top = Inches(1.8)

    panel = slide.shapes.add_textbox(Inches(0.8), top, Inches(12.0), Inches(5.4))
    tf = panel.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(21)


# Slide 1
add_title_slide(
    "Jira AI Copilot",
    "5-Minute Hackathon Demo\nReact UI → FastAPI → Jira RAG Engine"
)

# Slide 2: Idea
add_bullets_slide(
    "Idea",
    [
        "Turn Jira into a conversational decision assistant",
        "Ask natural questions; get evidence-backed answers instantly",
        "Focus on sprint risk, P1 ownership, and executive-ready summaries",
    ],
)

# Slide 3
add_bullets_slide(
    "Problem We Solved",
    [
        "Teams waste time searching Jira tickets, comments, and status history",
        "P1 ownership and sprint-risk visibility are often delayed",
        "Leaders need concise summaries with evidence, not guesswork",
    ],
)

# Slide 4: Tech stack
add_bullets_slide(
    "Tech Stack Used",
    [
        "Frontend: React (Vite)",
        "Backend: FastAPI (Python)",
        "AI/RAG: FAISS + embedding pipeline + retrieval filters",
        "Data source: Jira REST API (issues, comments, changelog)",
    ],
)

# Slide 5
add_bullets_slide(
    "Solution",
    [
        "Ask Jira questions in plain English",
        "RAG pipeline retrieves evidence from issues, comments, changelog",
        "Answers include structured summary, owners, and citations",
        "Supports executive summary + P1 intent (e.g., Any P1?)",
    ],
)

# Slide 6
add_demo_flow_slide()

# Slide 7: Code flow
add_text_panel_slide(
    "Code Flow",
    [
        "1) ingest_jira.py pulls Jira tickets/comments/changelog",
        "2) index_rag.py chunks data and builds FAISS index + metadata",
        "3) chatbot_cli.py retrieves top evidence and formats answer",
        "4) backend/api.py wraps chatbot logic as /ask endpoint",
        "5) frontend React calls /ask and shows answer + citations",
    ],
)

# Slide 8: Prompt use
add_text_panel_slide(
    "Prompt Use",
    [
        "Query example: 'Who is working on P1 tickets?'",
        "Intent detection: P1 phrases trigger priority filter",
        "Executive mode: appends 'executive summary' style handling",
        "Grounding rule: answer only from retrieved Jira evidence",
        "Output format: summary, owners, details, confidence, citations",
    ],
)

# Slide 9: Sequence flow
add_text_panel_slide(
    "Sequence Flow",
    [
        "User (UI) → FastAPI /ask",
        "FastAPI → load latest index bundle",
        "Retriever → FAISS search + metadata filters (project/P1/chunk)",
        "Context builder → assemble evidence blocks",
        "LLM/extractive formatter → structured response",
        "FastAPI → UI with citations and index used",
    ],
    subtitle="UI → API → Retrieval → Answer",
)

# Slide 10: Run commands
add_text_panel_slide(
    "Run Commands",
    [
        "1) Start stack: ./start_full_stack.sh",
        "2) Ingest data: python ingest_jira.py --site <tenant> --projects KAN --mode full",
        "3) Build index: python index_rag.py --snapshot-dir <snapshot> --provider hash",
        "4) Ask via CLI: python chatbot_cli.py --index-dir <index> --question 'Any P1?'",
        "5) Ask via UI: open http://localhost:5173 and click Ask Jira Copilot",
    ],
)

# Slide 11
add_highlight_slide(
    "What Makes This Trustworthy",
    "Evidence-backed AI for delivery decisions",
    [
        "Citations in every response",
        "Index used is shown in API/UI for traceability",
        "No-evidence behavior: explicit fallback instead of hallucination",
    ],
)

# Slide 12
add_bullets_slide(
    "Problems Faced & Fixes",
    [
        "Missing dependencies across Python environments → pinned and installed in pyenv",
        "Empty UI due to React import/runtime issue → fixed App import/build",
        "Stale citations from old index → auto-resolve latest index in backend",
        "No P1 detection initially → added intent-based priority filtering",
    ],
)

# Slide 13
add_bullets_slide(
    "Business Impact",
    [
        "Faster standups and incident triage",
        "Immediate visibility into P1 ownership",
        "Leadership-ready sprint summaries in seconds",
    ],
)

# Slide 14
add_bullets_slide(
    "Next Steps",
    [
        "Slack/Teams integration",
        "Role-based access and enterprise controls",
        "Trend analytics across sprints and projects",
    ],
)

prs.save(OUTPUT)
print(OUTPUT)
