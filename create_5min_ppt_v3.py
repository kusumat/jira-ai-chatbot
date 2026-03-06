from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT = "/Users/kusumathatavarthi/jira_ai_chatbot_artifacts/Jira_AI_Copilot_5min_Demo_v3.pptx"

# Personalization placeholders
PRESENTER_NAME = "Kusuma Thatavarthi"
TEAM_NAME = "Hackathon Team"
ORG_NAME = "Your Organization"

prs = Presentation()


def set_slide_bg_dark(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 14, 28)


def style_title(shape, color=RGBColor(230, 235, 255), size=42):
    tf = shape.text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.color.rgb = color
            run.font.size = Pt(size)
            run.font.bold = True


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg_dark(slide)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    style_title(slide.shapes.title, size=44)
    stf = slide.placeholders[1].text_frame
    for p in stf.paragraphs:
        for run in p.runs:
            run.font.color.rgb = RGBColor(145, 175, 255)
            run.font.size = Pt(20)


def add_dark_bullets_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_bg_dark(slide)
    slide.shapes.title.text = title
    style_title(slide.shapes.title, size=34)

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(220, 228, 255)


def add_conversation_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_dark(slide)
    slide.shapes.title.text = "Human + Robot: Delivery Briefing"
    style_title(slide.shapes.title, size=34)

    human_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.8), Inches(4.8))
    htf = human_box.text_frame
    htf.word_wrap = True
    human_lines = [
        "Human: Any P1 tickets?",
        "Human: Who owns them?",
        "Human: Will Sprint 24 slip?",
    ]
    for i, line in enumerate(human_lines):
        p = htf.paragraphs[0] if i == 0 else htf.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 218, 160)

    robot_box = slide.shapes.add_textbox(Inches(6.6), Inches(1.5), Inches(6.2), Inches(4.8))
    rtf = robot_box.text_frame
    rtf.word_wrap = True
    robot_lines = [
        "Robot: 5 P1 tickets detected.",
        "Robot: Owner is Kusuma Thatavarthi.",
        "Robot: Risk is elevated; CI and auth issues are driving delay.",
        "Robot: Evidence attached with Jira citations.",
    ]
    for i, line in enumerate(robot_lines):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = line
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(162, 250, 209)

    caption = slide.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(12), Inches(0.6))
    ctf = caption.text_frame
    ctf.text = "Original sci-fi inspired dialogue (copyright-safe)"
    ctf.paragraphs[0].font.size = Pt(14)
    ctf.paragraphs[0].font.color.rgb = RGBColor(129, 147, 191)


def add_demo_flow_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_dark(slide)
    slide.shapes.title.text = "Live Demo Flow (2 minutes)"
    style_title(slide.shapes.title, size=34)

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(12.0), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    steps = [
        "1) Ask: Any P1?",
        "2) Ask: Who is working on P1 tickets?",
        "3) Ask: Why was Sprint 24 delayed? executive summary",
        "4) Show citations + index used",
        "5) Close with business impact and next steps",
    ]
    for i, step in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = step
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(220, 228, 255)


def add_qa_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_dark(slide)
    slide.shapes.title.text = "Thank You · Q&A"
    style_title(slide.shapes.title, size=40)

    body = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.5), Inches(3.5))
    tf = body.text_frame
    tf.word_wrap = True

    lines = [
        f"Presenter: {PRESENTER_NAME}",
        f"Team: {TEAM_NAME}",
        f"Organization: {ORG_NAME}",
        "\nJira AI Copilot: Faster answers. Clear ownership. Trusted citations.",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(26 if i < 3 else 24)
        p.font.color.rgb = RGBColor(210, 223, 255)


add_title_slide(
    "Jira AI Copilot",
    f"5-Minute Demo · Sci-Fi Inspired Theme\nPresented by {PRESENTER_NAME} · {TEAM_NAME}",
)

add_dark_bullets_slide(
    "Problem",
    [
        "Jira context is fragmented across tickets, comments, and history",
        "P1 ownership takes too long to identify",
        "Leaders need concise, trusted sprint insights",
    ],
)

add_dark_bullets_slide(
    "Solution",
    [
        "Ask plain-language questions",
        "Retrieve Jira evidence with FAISS + metadata filters",
        "Return structured answer with citations and confidence",
    ],
)

add_conversation_slide()
add_demo_flow_slide()
add_dark_bullets_slide(
    "Business Impact",
    [
        "Faster standups with evidence-backed answers",
        "Immediate P1 ownership visibility",
        "Executive summaries in seconds",
        "Less Jira tab-hopping, more delivery focus",
    ],
)
add_dark_bullets_slide(
    "Roadmap",
    [
        "Slack/Teams bot integration",
        "Role-based controls and tenant governance",
        "Trend analytics across projects and sprints",
    ],
)
add_qa_slide()

prs.save(OUTPUT)
print(OUTPUT)
